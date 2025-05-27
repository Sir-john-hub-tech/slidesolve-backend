# main.py
from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import openai
import os
import sympy as sp
from sympy.parsing.sympy_parser import parse_expr
from typing import Dict, List
import json
import io
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI()

# CORS Configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize OpenAI
openai.api_key = os.getenv("OPENAI_API_KEY")

# In-memory storage with file validation
question_bank: Dict[str, Dict] = {}
student_answers: Dict[str, Dict] = {}
SUPPORTED_EXTENSIONS = {"pdf", "docx", "pptx", "ppt"}

# ====================== HELPER FUNCTIONS ======================
def extract_text(file_bytes: bytes, extension: str) -> str:
    """Extract text from supported file formats with error handling"""
    try:
        logger.info(f"Processing {extension.upper()} file")
        
        if extension == "pdf":
            with fitz.open(stream=file_bytes, filetype="pdf") as doc:
                return " ".join([page.get_text() for page in doc])
        
        elif extension in {"pptx", "ppt"}:
            prs = Presentation(io.BytesIO(file_bytes))
            return " ".join([
                shape.text 
                for slide in prs.slides 
                for shape in slide.shapes 
                if hasattr(shape, "text")
            ])
        
        elif extension == "docx":
            doc = Document(io.BytesIO(file_bytes))
            return " ".join([para.text for para in doc.paragraphs])
        
        raise ValueError(f"Unsupported file type: {extension}")
    
    except Exception as e:
        logger.error(f"Text extraction failed: {str(e)}")
        raise HTTPException(
            status_code=400,
            detail=f"Failed to process {extension.upper()} file: {str(e)}"
        )

def solve_math_problem(problem: str) -> Dict:
    """Solve mathematical equations with SymPy"""
    try:
        # Handle different equation formats
        problem = problem.replace("^", "**").strip()
        expr = parse_expr(problem)
        solution = sp.solve(expr)
        steps = sp.pretty(solution)
        return {
            "problem": problem,
            "solution": str(solution),
            "steps": steps
        }
    except Exception as e:
        logger.error(f"Math solving failed: {str(e)}")
        return {
            "error": f"Could not solve {problem}",
            "details": str(e)
        }

def generate_exam_questions(text: str) -> Dict:
    """Generate 50 questions using OpenAI with validation"""
    try:
        prompt = f"""Generate 50 comprehensive exam questions from this text.
        Include multiple choice, fill-in-the-blank, and short answer questions.
        Use this exact JSON format:
        {{
            "multiple_choice": [
                {{"question": "...", "options": ["A","B","C","D"], "answer": "A"}}
            ],
            "fill_in": [
                {{"question": "...", "answer": "..."}}
            ],
            "short_answer": [
                {{"question": "...", "answer": "..."}}
            ]
        }}
        Text content: {text[:3000]}"""  # Limit text length for OpenAI

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{
                "role": "user",
                "content": prompt
            }],
            temperature=0.3  # More deterministic output
        )
        
        # Validate and parse response
        content = response.choices[0].message.content
        try:
            return json.loads(content)
        except json.JSONDecodeError:
            logger.error("Invalid JSON from OpenAI: %s", content)
            raise HTTPException(
                status_code=500,
                detail="Failed to parse question format. Please try again."
            )
    
    except openai.error.AuthenticationError:
        raise HTTPException(
            status_code=401,
            detail="Invalid OpenAI API key. Check your configuration."
        )
    except Exception as e:
        logger.error(f"OpenAI API Error: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Question generation failed: {str(e)}"
        )

# ====================== API ENDPOINTS ======================
@app.get("/")
async def welcome():
    """Display welcome message with emoji"""
    return {"message": "Hello student 😊, welcome to Sir John's learning tool. Enjoy!"}

@app.post("/upload/")
async def upload_file(file: UploadFile = File(...)):
    """Process uploaded files and generate questions"""
    try:
        # Validate file type
        extension = file.filename.split(".")[-1].lower()
        if extension not in SUPPORTED_EXTENSIONS:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type. Supported formats: {SUPPORTED_EXTENSIONS}"
            )
        
        # Process file
        file_bytes = await file.read()
        text = extract_text(file_bytes, extension)
        questions = generate_exam_questions(text)
        
        # Store results
        question_bank[file.filename] = {
            "text": text,
            "questions": questions
        }
        
        return {
            "filename": file.filename,
            "questions": questions,
            "message": "File processed successfully"
        }
    
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Unexpected error: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Internal server error: {str(e)}"
        )

@app.post("/solve-math/")
async def math_solver(problem: str = Form(...)):
    """Solve mathematical equations"""
    result = solve_math_problem(problem)
    if "error" in result:
        raise HTTPException(
            status_code=400,
            detail=result["error"]
        )
    return result

@app.post("/submit-answers/")
async def submit_answers(filename: str = Form(...), answers: str = Form(...)):
    """Store and grade student answers"""
    try:
        # Validate answers format
        try:
            answers_dict = json.loads(answers)
        except json.JSONDecodeError:
            raise HTTPException(
                status_code=400,
                detail="Invalid answer format. Use JSON."
            )
        
        # Store answers
        student_answers[filename] = answers_dict
        return {"message": "Answers submitted successfully"}
    
    except Exception as e:
        logger.error(f"Answer submission failed: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process answers: {str(e)}"
        )

@app.get("/results/{filename}")
async def get_results(filename: str):
    """Provide graded results and study suggestions"""
    try:
        # Get stored data
        questions = question_bank.get(filename, {})
        answers = student_answers.get(filename, {})
        
        # Calculate score
        correct = 0
        total = 0
        feedback = []
        
        for q_type in ["multiple_choice", "fill_in", "short_answer"]:
            for question in questions.get("questions", {}).get(q_type, []):
                total += 1
                user_answer = answers.get(question["question"], "").strip().lower()
                correct_answer = question["answer"].strip().lower()
                
                if user_answer == correct_answer:
                    correct += 1
                else:
                    feedback.append({
                        "question": question["question"],
                        "your_answer": user_answer,
                        "correct_answer": correct_answer
                    })
        
        # Generate suggestions
        score = (correct / total) * 100 if total > 0 else 0
        suggestions = []
        
        if score < 50:
            suggestions = [
                "Focus on fundamental concepts",
                "Review chapter summaries",
                "Practice basic definitions"
            ]
        elif score < 75:
            suggestions = [
                "Work on application problems",
                "Practice time management",
                "Review diagrams and charts"
            ]
        else:
            suggestions = [
                "Excellent performance!",
                "Challenge yourself with advanced problems",
                "Help peers with difficult concepts"
            ]
        
        return {
            "score": f"{score:.1f}%",
            "correct": correct,
            "total": total,
            "feedback": feedback[:5],  # Show top 5 mistakes
            "suggestions": suggestions
        }
    
    except Exception as e:
        logger.error(f"Result calculation failed: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate results: {str(e)}"
        )

# ====================== REQUIREMENTS.TXT ======================
"""
fastapi==0.103.1
uvicorn==0.23.2
python-multipart==0.0.6
pymupdf==1.23.5
python-pptx==0.6.23
python-docx==1.1.0
openai==1.3.6
sympy==1.12
python-dotenv==1.0.0
"""
